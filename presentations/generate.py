#!/usr/bin/env python3
"""Generate a branded PowerPoint (.pptx) from a markdown input file.

This script parses a simple markdown document (title, optional blockquote
subtitle, and level-2 headings with bullet lists) and renders it as a
branded slide deck: a navy title slide, one content slide per section, and
a closing slide. After building, it appends a record to a memory log.
"""

import argparse
import datetime
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand colors
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
GRAY = RGBColor(0x5A, 0x63, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def parse_markdown(text):
    """Parse markdown text into a structured dict.

    Returns: {"title": str, "subtitle": str,
              "slides": [{"heading": str, "bullets": [(level, text), ...]}]}
    """
    title = ""
    subtitle = ""
    slides = []
    current = None  # current content slide being built

    # Lines that are not a title and not empty, used for the fallback "Notes"
    # slide when there are no level-2 sections.
    misc_lines = []

    lines = text.splitlines()

    for raw in lines:
        line = raw.rstrip("\n")
        stripped = line.strip()

        # Skip blank lines.
        if not stripped:
            continue

        # Title: first level-1 heading.
        if stripped.startswith("# ") and not title:
            title = stripped[2:].strip()
            continue

        # Subtitle: first blockquote line in the document.
        if stripped.startswith("> ") and not subtitle:
            subtitle = stripped[2:].strip()
            continue

        # Content slide: level-2 heading.
        if stripped.startswith("## "):
            current = {"heading": stripped[3:].strip(), "bullets": []}
            slides.append(current)
            continue

        # Bullets: lines starting with "- " or "* ", possibly indented.
        bullet_match = re.match(r"^(\s*)([-*])\s+(.*)$", line)
        if bullet_match:
            indent = bullet_match.group(1)
            bullet_text = bullet_match.group(3).strip()
            # One level of indentation (2 or 4 leading spaces) -> sub-bullet.
            level = 1 if len(indent) >= 2 else 0
            if current is not None:
                current["bullets"].append((level, bullet_text))
            else:
                # No section yet; remember for fallback "Notes" slide.
                misc_lines.append(bullet_text)
            continue

        # Any other non-empty, non-title line: keep for fallback slide.
        misc_lines.append(stripped)

    # Default subtitle if none was found.
    if not subtitle:
        today = datetime.date.today().strftime("%B %-d, %Y")
        subtitle = "Weekly Update — " + today

    # If there were no level-2 sections, build a single "Notes" slide.
    if not slides:
        bullets = [(0, ln) for ln in misc_lines]
        slides.append({"heading": "Notes", "bullets": bullets})

    return {"title": title or "Untitled", "subtitle": subtitle, "slides": slides}


def _set_slide_background(slide, color):
    """Set a slide's background fill to a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle):
    """Add the navy title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, NAVY)

    # Title textbox.
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.6)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Subtitle textbox.
    sub_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.9), Inches(8.4), Inches(1.0)
    )
    stf = sub_box.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    srun = sp.add_run()
    srun.text = subtitle
    srun.font.size = Pt(20)
    srun.font.color.rgb = ACCENT

    # Footer textbox.
    footer_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(6.9), Inches(8.4), Inches(0.4)
    )
    ftf = footer_box.text_frame
    fp = ftf.paragraphs[0]
    frun = fp.add_run()
    frun.text = "Momentum 360 · Weekly Presentation"
    frun.font.size = Pt(12)
    frun.font.color.rgb = GRAY


def add_content_slide(prs, heading, bullets):
    """Add a white content slide with an accent bar and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, WHITE)

    # Accent bar across the top.
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.25)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Heading textbox.
    head_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.5), Inches(8.8), Inches(1.0)
    )
    htf = head_box.text_frame
    htf.word_wrap = True
    hp = htf.paragraphs[0]
    hrun = hp.add_run()
    hrun.text = heading
    hrun.font.size = Pt(28)
    hrun.font.bold = True
    hrun.font.color.rgb = NAVY

    # Body textbox with bullets.
    body_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.7), Inches(8.4), Inches(4.8)
    )
    btf = body_box.text_frame
    btf.word_wrap = True

    first = True
    for level, text in bullets:
        # Reuse the empty first paragraph, otherwise add a new one.
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        if level == 0:
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = "•  " + text
            run.font.size = Pt(18)
            run.font.color.rgb = NAVY
        else:
            p.level = 1
            run = p.add_run()
            run.text = "–  " + text
            run.font.size = Pt(15)
            run.font.color.rgb = GRAY


def add_closing_slide(prs):
    """Add the navy closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, NAVY)

    box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.8)
    )
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank you"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    today = datetime.date.today().strftime("%B %-d, %Y")
    run2 = p2.add_run()
    run2.text = "Generated " + today
    run2.font.size = Pt(16)
    run2.font.color.rgb = ACCENT


def build(input_path, output_path):
    """Read markdown, build the deck, and save it. Returns the parsed dict."""
    with open(input_path, "r", encoding="utf-8") as fh:
        text = fh.read()

    parsed = parse_markdown(text)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, parsed["title"], parsed["subtitle"])
    for slide in parsed["slides"]:
        add_content_slide(prs, slide["heading"], slide["bullets"])
    add_closing_slide(prs)

    prs.save(output_path)
    return parsed


def _resolve(path, repo_root):
    """Resolve a path under the repo root if it is not absolute."""
    p = Path(path)
    if not p.is_absolute():
        p = repo_root / p
    return p


def main():
    repo_root = Path(__file__).resolve().parent.parent

    today = datetime.date.today().strftime("%Y-%m-%d")
    default_output = "presentations/output/%s-weekly.pptx" % today

    parser = argparse.ArgumentParser(
        description="Generate a branded PowerPoint from a markdown file."
    )
    parser.add_argument("--input", default="presentations/input/next.md")
    parser.add_argument("--output", default=default_output)
    parser.add_argument("--memory", default="presentations/MEMORY.md")
    args = parser.parse_args()

    input_path = _resolve(args.input, repo_root)
    output_path = _resolve(args.output, repo_root)
    memory_path = _resolve(args.memory, repo_root)

    # Ensure the output directory exists.
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if not input_path.exists():
        sys.stderr.write("Error: input file not found: %s\n" % input_path)
        sys.exit(1)

    parsed = build(str(input_path), str(output_path))
    n_slides = len(parsed["slides"])

    # Append a record to the memory log.
    now = datetime.datetime.now(datetime.timezone.utc)
    timestamp = now.strftime("%Y-%m-%d %H:%M UTC")
    try:
        rel_output = output_path.relative_to(repo_root)
    except ValueError:
        rel_output = output_path
    memory_line = '- %s — Generated "%s" (%d content slides) → %s\n' % (
        timestamp,
        parsed["title"],
        n_slides,
        rel_output,
    )
    memory_path.parent.mkdir(parents=True, exist_ok=True)
    with open(memory_path, "a", encoding="utf-8") as fh:
        fh.write(memory_line)

    # Success summary.
    print("Generated presentation: %s" % output_path)
    print('Title: "%s"' % parsed["title"])
    print("Content slides: %d" % n_slides)


if __name__ == "__main__":
    main()
